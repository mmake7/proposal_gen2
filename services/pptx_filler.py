"""
PPTX Filler
템플릿 PPTX에 생성된 콘텐츠를 채우고 폰트 프리셋을 적용합니다.

매핑된 챕터의 슬라이드에 TOC 제목과 생성된 콘텐츠를 삽입하고,
폰트 설정 모드에 따라 폰트를 적용하여 최종 PPTX를 생성합니다.
"""

from __future__ import annotations

import io
import logging
import os

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor

from services.font_manager import get_font_settings, get_preset

logger = logging.getLogger(__name__)

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATES_STORAGE_DIR = os.path.join(BASE_DIR, 'data', 'templates_storage')


class PptxFillerError(Exception):
    """PPTX 채우기 중 발생하는 에러"""


def fill_template(
    template_id: str,
    mapping: list[dict],
    content: dict,
    toc_items: list[dict],
) -> io.BytesIO:
    """템플릿 PPTX에 콘텐츠를 채워 새 파일을 생성합니다.

    Args:
        template_id: 템플릿 ID
        mapping: TOC↔챕터 매핑 리스트
        content: 생성된 콘텐츠 {sections: {"0-0": {...}, ...}}
        toc_items: 확정된 TOC 항목 리스트

    Returns:
        BytesIO containing the filled PPTX file
    """
    # 원본 템플릿 로드
    original_path = os.path.join(TEMPLATES_STORAGE_DIR, template_id, 'template.pptx')
    if not os.path.exists(original_path):
        raise PptxFillerError('원본 템플릿 파일을 찾을 수 없습니다.')

    prs = Presentation(original_path)

    # 폰트 설정 로드
    font_settings = get_font_settings()
    font_mode = font_settings.get('mode', 'template_keep')
    preset = None
    if font_mode in ('preset_apply', 'preset_force'):
        preset_id = font_settings.get('active_preset_id')
        if preset_id:
            preset = get_preset(preset_id)

    sections = content.get('sections', {})

    # 매핑된 챕터별로 콘텐츠 채우기
    for entry in mapping:
        if entry.get('status') == 'unmatched':
            continue

        toc_idx = entry.get('toc_index', 0)
        toc_title = entry.get('toc_title', '')
        slide_numbers = entry.get('slides', [])

        if not slide_numbers:
            continue

        # 해당 TOC L1의 L2 자식들의 콘텐츠 수집
        l1_item = toc_items[toc_idx] if toc_idx < len(toc_items) else None
        l2_children = (l1_item.get('children', []) if l1_item else [])
        section_contents = []

        if l2_children:
            for l2_idx in range(len(l2_children)):
                key = f'{toc_idx}-{l2_idx}'
                sec = sections.get(key)
                if sec:
                    section_contents.append(sec)
        else:
            # L1에 자식이 없으면 L1 자체의 콘텐츠
            key = str(toc_idx)
            sec = sections.get(key)
            if sec:
                section_contents.append(sec)

        # 슬라이드에 콘텐츠 채우기
        _fill_chapter_slides(prs, slide_numbers, toc_title, section_contents, font_mode, preset)

    # preset_force 모드: 전체 슬라이드에 폰트 적용
    if font_mode == 'preset_force' and preset:
        _apply_font_to_all(prs, preset)

    # BytesIO로 저장
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _fill_chapter_slides(
    prs: Presentation,
    slide_numbers: list[int],
    chapter_title: str,
    section_contents: list[dict],
    font_mode: str,
    preset: dict | None,
):
    """챕터 슬라이드에 콘텐츠를 채웁니다."""
    slides = list(prs.slides)

    for i, slide_num in enumerate(slide_numbers):
        slide_idx = slide_num - 1  # 0-indexed
        if slide_idx < 0 or slide_idx >= len(slides):
            continue

        slide = slides[slide_idx]

        if i == 0:
            # 첫 슬라이드: 챕터 제목 슬라이드
            _fill_title_slide(slide, chapter_title, font_mode, preset)
        else:
            # 나머지 슬라이드: 콘텐츠 채우기
            content_idx = i - 1
            if content_idx < len(section_contents):
                _fill_content_slide(
                    slide, section_contents[content_idx],
                    font_mode, preset,
                )


def _fill_title_slide(
    slide,
    title: str,
    font_mode: str,
    preset: dict | None,
):
    """제목 슬라이드에 챕터 제목을 채웁니다."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # 제목 placeholder 또는 가장 큰 텍스트 영역 찾기
        if shape.shape_id == 0 or _is_title_placeholder(shape):
            tf = shape.text_frame
            if tf.paragraphs:
                para = tf.paragraphs[0]
                # 기존 런 제거 후 새 텍스트 설정
                for run in para.runs:
                    run.text = ''
                if para.runs:
                    para.runs[0].text = title
                    if font_mode == 'preset_apply' and preset:
                        _apply_font_to_run(para.runs[0], preset.get('title_font'))
                else:
                    run = para.add_run()
                    run.text = title
                    if font_mode == 'preset_apply' and preset:
                        _apply_font_to_run(run, preset.get('title_font'))
            break


def _fill_content_slide(
    slide,
    content: dict,
    font_mode: str,
    preset: dict | None,
):
    """콘텐츠 슬라이드에 텍스트를 채웁니다."""
    title_text = content.get('title', '')
    body_text = content.get('content_text', '')
    bullet_points = content.get('bullet_points', [])

    title_filled = False
    body_filled = False

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        if not title_filled and _is_title_placeholder(shape):
            _set_text_frame(shape.text_frame, title_text, font_mode, preset, 'title_font')
            title_filled = True
        elif not body_filled and _is_body_placeholder(shape):
            # 본문: 텍스트 + 불릿 포인트
            lines = []
            if body_text:
                lines.append(body_text)
            for bp in bullet_points:
                lines.append(f'• {bp}')
            full_text = '\n\n'.join(lines) if lines else ''
            _set_text_frame(shape.text_frame, full_text, font_mode, preset, 'body_font')
            body_filled = True

        if title_filled and body_filled:
            break


def _is_title_placeholder(shape) -> bool:
    """shape가 제목 placeholder인지 확인합니다."""
    try:
        ph = shape.placeholder_format
        if ph is not None:
            return ph.idx in (0, 1)  # Title or Center Title
    except Exception:
        pass
    return False


def _is_body_placeholder(shape) -> bool:
    """shape가 본문 placeholder인지 확인합니다."""
    try:
        ph = shape.placeholder_format
        if ph is not None:
            return ph.idx in (1, 2, 13, 14)  # Body, Subtitle, etc.
    except Exception:
        pass
    # 큰 텍스트 프레임을 본문으로 간주
    if shape.has_text_frame:
        try:
            width = shape.width
            if width and int(width) > 3000000:  # > ~2.3 inches
                return True
        except Exception:
            pass
    return False


def _set_text_frame(
    text_frame,
    text: str,
    font_mode: str,
    preset: dict | None,
    font_key: str,
):
    """텍스트 프레임에 텍스트를 설정합니다."""
    if not text:
        return

    # 기존 텍스트 제거
    for para in text_frame.paragraphs:
        for run in para.runs:
            run.text = ''

    # 첫 단락에 텍스트 설정
    paragraphs = text.split('\n\n')
    for p_idx, p_text in enumerate(paragraphs):
        if p_idx == 0 and text_frame.paragraphs:
            para = text_frame.paragraphs[0]
        else:
            para = text_frame.add_paragraph()

        if para.runs:
            para.runs[0].text = p_text.strip()
            if font_mode == 'preset_apply' and preset:
                _apply_font_to_run(para.runs[0], preset.get(font_key))
        else:
            run = para.add_run()
            run.text = p_text.strip()
            if font_mode == 'preset_apply' and preset:
                _apply_font_to_run(run, preset.get(font_key))


def _apply_font_to_run(run, font_def: dict | None):
    """런에 폰트 프리셋을 적용합니다."""
    if not font_def:
        return

    font = run.font

    name = font_def.get('name')
    if name:
        font.name = name

    size_pt = font_def.get('size_pt')
    if size_pt:
        font.size = Pt(float(size_pt))

    bold = font_def.get('bold')
    if bold is not None:
        font.bold = bool(bold)

    color = font_def.get('color')
    if color and isinstance(color, str) and color.startswith('#'):
        hex_str = color.lstrip('#')
        if len(hex_str) == 6:
            try:
                font.color.rgb = RGBColor(
                    int(hex_str[0:2], 16),
                    int(hex_str[2:4], 16),
                    int(hex_str[4:6], 16),
                )
            except (ValueError, TypeError):
                pass


def _apply_font_to_all(prs: Presentation, preset: dict):
    """전체 프레젠테이션에 프리셋 폰트를 적용합니다 (preset_force 모드)."""
    title_font = preset.get('title_font')
    body_font = preset.get('body_font')

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            is_title = _is_title_placeholder(shape)

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not run.text.strip():
                        continue
                    if is_title:
                        _apply_font_to_run(run, title_font)
                    else:
                        _apply_font_to_run(run, body_font)
