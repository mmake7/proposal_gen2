"""
PowerPoint Template Manager
업로드된 PPTX 파일의 슬라이드 구조를 분석하고 장(Chapter) 단위로 그룹핑합니다.

핵심 기능:
- 장 구분 슬라이드 자동 감지 (제목만 있는 슬라이드: I.제안업체 일반, II.전략 등)
- 각 슬라이드별 플레이스홀더 목록 (위치, 크기, 타입, 인덱스)
- 장별 소속 슬라이드 그룹핑
- 슬라이드 썸네일 PNG 생성
- 템플릿 메타정보 JSON 관리
"""

from __future__ import annotations

import io
import json
import logging
import os
import re
import uuid
from datetime import datetime, timezone

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)

# 저장 경로
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATES_STORAGE_DIR = os.path.join(BASE_DIR, 'data', 'templates_storage')
METADATA_FILE = os.path.join(TEMPLATES_STORAGE_DIR, '_metadata.json')

# 장 구분 패턴 (로마숫자/아라비아숫자 + 점/마침표 + 제목)
CHAPTER_PATTERNS = [
    re.compile(r'^[IVXivx]+[\.\s]'),          # I. 제안업체 일반 / IV. 프로젝트 관리
    re.compile(r'^[IVXivx]+\s*[\.\-\:]\s*'),   # I - 제안업체 / I: 전략
    re.compile(r'^[ⅠⅡⅢⅣⅤⅥⅦⅧⅨⅩⅪⅫ][\.\s\-\:]'),  # Ⅰ. 제안업체 / Ⅱ. 전략 (전각 로마숫자)
    re.compile(r'^[ⅰⅱⅲⅳⅴⅵⅶⅷⅸⅹⅺⅻ][\.\s\-\:]'),  # ⅰ. / ⅱ. (소문자 전각 로마숫자)
    re.compile(r'^\d+[\.\s]'),                  # 1. 제안업체 일반
    re.compile(r'^제?\d+장[\.\s:]'),             # 제1장. 제안개요 / 1장.
    re.compile(r'^Part\s*\d+', re.IGNORECASE),  # Part 1 / PART 2
]


class TemplateManagerError(Exception):
    """템플릿 관리 중 발생하는 에러"""


# ── Storage Helpers ──────────────────────────────────────────

def _ensure_storage_dir():
    os.makedirs(TEMPLATES_STORAGE_DIR, exist_ok=True)


def _load_metadata() -> list[dict]:
    _ensure_storage_dir()
    if not os.path.exists(METADATA_FILE):
        return []
    with open(METADATA_FILE, 'r', encoding='utf-8') as f:
        return json.load(f)


def _save_metadata(metadata: list[dict]):
    _ensure_storage_dir()
    with open(METADATA_FILE, 'w', encoding='utf-8') as f:
        json.dump(metadata, f, ensure_ascii=False, indent=2)


def _template_dir(template_id: str) -> str:
    """개별 템플릿 저장 디렉터리"""
    d = os.path.join(TEMPLATES_STORAGE_DIR, template_id)
    os.makedirs(d, exist_ok=True)
    return d


# ── Unit Conversion ──────────────────────────────────────────

def _emu_to_cm(emu_val) -> float | None:
    if emu_val is None:
        return None
    return round(int(emu_val) / 914400 * 2.54, 2)


def _emu_to_pt(emu_val) -> float | None:
    if emu_val is None:
        return None
    return round(int(emu_val) / 12700, 1)


# ── Chapter Detection ────────────────────────────────────────

def _is_chapter_divider(slide) -> bool:
    """장 구분 슬라이드인지 판별합니다.

    장 구분 슬라이드 특징:
    1. 텍스트를 포함한 셰이프가 1~2개뿐 (제목만 있음)
    2. 제목 텍스트가 장 구분 패턴(I., II., 1., 제1장 등)과 일치
    3. 본문 콘텐츠가 거의 없음
    """
    text_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                text_shapes.append(text)

    # 텍스트 셰이프가 0개면 구분 슬라이드 아님 (빈 슬라이드)
    if len(text_shapes) == 0:
        return False

    # 텍스트 셰이프가 3개 이상이면 일반 콘텐츠 슬라이드
    if len(text_shapes) > 2:
        return False

    # 주요 텍스트(첫 번째)가 장 구분 패턴과 일치하는지 확인
    primary_text = text_shapes[0]

    # 패턴 매칭
    for pattern in CHAPTER_PATTERNS:
        if pattern.search(primary_text):
            return True

    # 레이아웃 이름으로 추가 판별
    if slide.slide_layout:
        layout_lower = (slide.slide_layout.name or '').lower()
        if any(keyword in layout_lower for keyword in
               ['section', 'divider', '구분', '장표지', 'title only', '제목만']):
            return True

    return False


def _extract_chapter_title(slide) -> str:
    """장 구분 슬라이드에서 제목을 추출합니다."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                return text
    return '(제목 없음)'


# ── Placeholder Analysis ─────────────────────────────────────

def _analyze_placeholder(ph) -> dict:
    """단일 플레이스홀더의 상세 정보를 추출합니다."""
    ph_type = str(ph.placeholder_format.type) if ph.placeholder_format.type else 'UNKNOWN'

    return {
        'idx': ph.placeholder_format.idx,
        'type': ph_type,
        'name': ph.name or '',
        'left_cm': _emu_to_cm(ph.left),
        'top_cm': _emu_to_cm(ph.top),
        'width_cm': _emu_to_cm(ph.width),
        'height_cm': _emu_to_cm(ph.height),
        'text': (ph.text or '')[:500],
    }


def _analyze_shapes(slide) -> list[dict]:
    """슬라이드 내 모든 셰이프의 요약 정보를 추출합니다."""
    shapes = []
    for shape in slide.shapes:
        info = {
            'name': shape.name or '',
            'shape_type': str(shape.shape_type) if shape.shape_type else 'UNKNOWN',
            'left_cm': _emu_to_cm(shape.left),
            'top_cm': _emu_to_cm(shape.top),
            'width_cm': _emu_to_cm(shape.width),
            'height_cm': _emu_to_cm(shape.height),
        }
        if shape.has_text_frame:
            info['text'] = shape.text_frame.text.strip()[:300]
        shapes.append(info)
    return shapes


# ── Font Extraction ──────────────────────────────────────────

def _extract_font_info(slide) -> list[dict]:
    """슬라이드에서 사용된 폰트 정보를 추출합니다.

    python-pptx run.font 객체에서 접근:
    - run.font.name: 폰트 이름 (None이면 테마 폰트)
    - run.font.size: 폰트 크기 (EMU, 1pt = 12700 EMU)
    - run.font.bold: 볼드 여부 (True/False/None)
    - run.font.color.rgb: RGBColor 객체

    Returns:
        슬라이드별 고유 폰트 정보 리스트
    """
    seen = set()
    font_entries = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if not run.text.strip():
                    continue

                font = run.font
                font_name = font.name

                # 크기: EMU → pt
                size_pt = None
                if font.size is not None:
                    size_pt = round(int(font.size) / 12700, 1)

                # 볼드
                is_bold = bool(font.bold) if font.bold is not None else False

                # 색상
                color_hex = None
                try:
                    if font.color and font.color.type is not None:
                        rgb = font.color.rgb
                        if rgb is not None:
                            color_hex = '#' + str(rgb)
                except Exception:
                    pass

                if font_name is None and size_pt is None:
                    continue

                key = (font_name, size_pt, is_bold, color_hex)
                if key in seen:
                    continue
                seen.add(key)

                font_entries.append({
                    'font_name': font_name or '(테마 폰트)',
                    'size_pt': size_pt,
                    'is_bold': is_bold,
                    'color': color_hex,
                })

    return font_entries


# ── Main Analysis ────────────────────────────────────────────

def analyze_pptx(file_bytes: bytes, filename: str) -> dict:
    """PPTX 파일의 슬라이드 구조를 분석하고 장별로 그룹핑합니다.

    Returns:
        {
            filename, slide_count, slide_width_cm, slide_height_cm,
            chapters: [
                {
                    title: "I. 제안업체 일반",
                    divider_slide_index: 1,
                    slides: [
                        {index, layout_name, is_chapter_divider, placeholders, shapes, text_content},
                        ...
                    ]
                },
                ...
            ],
            all_slides: [ ... ]  # 전체 슬라이드 flat 목록
        }
    """
    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception as exc:
        raise TemplateManagerError(f'PPTX 파일을 열 수 없습니다: {exc}') from exc

    slide_width_cm = _emu_to_cm(prs.slide_width)
    slide_height_cm = _emu_to_cm(prs.slide_height)

    all_slides = []
    chapters = []
    current_chapter = None

    for idx, slide in enumerate(prs.slides):
        layout_name = ''
        if slide.slide_layout:
            layout_name = slide.slide_layout.name or ''

        is_divider = _is_chapter_divider(slide)

        # 플레이스홀더 분석
        placeholders = [_analyze_placeholder(ph) for ph in slide.placeholders]

        # 셰이프 분석
        shapes = _analyze_shapes(slide)

        # 텍스트 추출
        text_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        text_parts.append(t)

        # 폰트 정보 추출
        fonts = _extract_font_info(slide)

        slide_info = {
            'index': idx + 1,
            'layout_name': layout_name,
            'is_chapter_divider': is_divider,
            'placeholders': placeholders,
            'shapes': shapes,
            'text_content': '\n'.join(text_parts)[:2000],
            'fonts': fonts,
        }

        all_slides.append(slide_info)

        # 장 구분 슬라이드를 만나면 새 챕터 시작
        if is_divider:
            chapter_title = _extract_chapter_title(slide)
            current_chapter = {
                'title': chapter_title,
                'divider_slide_index': idx + 1,
                'slides': [slide_info],
            }
            chapters.append(current_chapter)
        else:
            if current_chapter is None:
                # 첫 장 구분 전의 슬라이드들 → "표지/서문" 챕터
                current_chapter = {
                    'title': '표지 / 서문',
                    'divider_slide_index': None,
                    'slides': [],
                }
                chapters.append(current_chapter)
            current_chapter['slides'].append(slide_info)

    # 빈 챕터 없으면 전체를 하나의 챕터로
    if not chapters:
        chapters = [{
            'title': '전체',
            'divider_slide_index': None,
            'slides': all_slides,
        }]

    # 폰트 요약 집계
    font_usage: dict[str, dict] = {}
    for s in all_slides:
        for fi in s.get('fonts', []):
            name = fi['font_name']
            if name not in font_usage:
                font_usage[name] = {'count': 0, 'sizes': set(), 'has_bold': False}
            font_usage[name]['count'] += 1
            if fi['size_pt'] is not None:
                font_usage[name]['sizes'].add(fi['size_pt'])
            if fi['is_bold']:
                font_usage[name]['has_bold'] = True

    font_summary = {
        'unique_fonts': sorted(font_usage.keys()),
        'font_count': len(font_usage),
        'details': [
            {
                'font_name': name,
                'usage_count': info['count'],
                'sizes_used': sorted(info['sizes']),
                'has_bold': info['has_bold'],
            }
            for name, info in sorted(font_usage.items())
        ],
    }

    return {
        'filename': filename,
        'slide_count': len(all_slides),
        'slide_width_cm': slide_width_cm,
        'slide_height_cm': slide_height_cm,
        'chapter_count': len(chapters),
        'chapters': chapters,
        'all_slides': all_slides,
        'font_summary': font_summary,
    }


# ── Thumbnail Generation ─────────────────────────────────────

def generate_slide_thumbnail(file_bytes: bytes, slide_num: int,
                             width: int = 480) -> bytes | None:
    """특정 슬라이드의 썸네일 PNG를 생성합니다.

    python-pptx 자체에는 렌더링 기능이 없으므로,
    슬라이드의 셰이프 정보를 기반으로 Pillow로 단순한 레이아웃 프리뷰를 생성합니다.
    실제 PowerPoint 렌더링이 아닌, 플레이스홀더 위치를 시각화한 와이어프레임입니다.

    Args:
        file_bytes: PPTX 파일 바이트
        slide_num: 슬라이드 번호 (1-based)
        width: 출력 이미지 너비 (px)

    Returns:
        PNG 이미지 바이트, 실패시 None
    """
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        logger.warning('Pillow가 설치되지 않아 썸네일을 생성할 수 없습니다.')
        return None

    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception:
        return None

    if slide_num < 1 or slide_num > len(prs.slides):
        return None

    slide = prs.slides[slide_num - 1]

    # 슬라이드 크기 → 이미지 크기 계산
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    scale = width / slide_w
    img_h = int(slide_h * scale)

    # 흰색 배경 이미지
    img = Image.new('RGB', (width, img_h), '#FFFFFF')
    draw = ImageDraw.Draw(img)

    # 슬라이드 배경 테두리
    draw.rectangle([0, 0, width - 1, img_h - 1], outline='#E2E8F0', width=2)

    # 폰트 (시스템 기본 폰트 사용)
    try:
        font_small = ImageFont.truetype('arial.ttf', max(10, int(12 * scale * 2)))
        font_title = ImageFont.truetype('arial.ttf', max(12, int(18 * scale * 2)))
    except (IOError, OSError):
        font_small = ImageFont.load_default()
        font_title = font_small

    # 셰이프 그리기
    for shape in slide.shapes:
        x = int(shape.left * scale)
        y = int(shape.top * scale)
        w = int(shape.width * scale)
        h = int(shape.height * scale)

        # 범위 보정
        x = max(0, min(x, width - 2))
        y = max(0, min(y, img_h - 2))
        w = max(4, min(w, width - x))
        h = max(4, min(h, img_h - y))

        is_placeholder = shape.is_placeholder if hasattr(shape, 'is_placeholder') else False

        if is_placeholder:
            # 플레이스홀더: 연한 인디고 배경 + 테두리
            draw.rectangle([x, y, x + w, y + h],
                           fill='#EEF2FF', outline='#A5B4FC', width=1)
        else:
            # 일반 셰이프: 연한 회색 배경
            draw.rectangle([x, y, x + w, y + h],
                           fill='#F8FAFC', outline='#E2E8F0', width=1)

        # 텍스트 라벨
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                label = text[:30] + ('...' if len(text) > 30 else '')
                font = font_title if is_placeholder else font_small
                text_x = x + 4
                text_y = y + 4
                # 텍스트가 영역 안에 들어가는지 확인
                if text_x + 10 < x + w and text_y + 10 < y + h:
                    draw.text((text_x, text_y), label,
                              fill='#475569', font=font)

    # PNG로 인코딩
    buf = io.BytesIO()
    img.save(buf, format='PNG', optimize=True)
    buf.seek(0)
    return buf.getvalue()


# ── CRUD Operations ──────────────────────────────────────────

def save_template(file_bytes: bytes, filename: str) -> dict:
    """PPTX 파일을 분석 후 저장하고 메타정보를 반환합니다."""
    # 1. 분석
    analysis = analyze_pptx(file_bytes, filename)

    # 2. ID 생성
    template_id = str(uuid.uuid4())[:8]

    # 3. 파일 저장
    tpl_dir = _template_dir(template_id)
    pptx_path = os.path.join(tpl_dir, 'template.pptx')
    with open(pptx_path, 'wb') as f:
        f.write(file_bytes)

    # 4. 분석 결과 저장
    analysis_path = os.path.join(tpl_dir, 'analysis.json')
    with open(analysis_path, 'w', encoding='utf-8') as f:
        json.dump(analysis, f, ensure_ascii=False, indent=2)

    # 5. 메타정보
    now = datetime.now(timezone.utc).isoformat()
    font_summary = analysis.get('font_summary', {})
    meta_entry = {
        'id': template_id,
        'name': os.path.splitext(filename)[0],
        'original_filename': filename,
        'slide_count': analysis['slide_count'],
        'chapter_count': analysis['chapter_count'],
        'chapter_titles': [ch['title'] for ch in analysis['chapters']],
        'slide_width_cm': analysis['slide_width_cm'],
        'slide_height_cm': analysis['slide_height_cm'],
        'font_count': font_summary.get('font_count', 0),
        'fonts_used': font_summary.get('unique_fonts', []),
        'uploaded_at': now,
    }

    metadata = _load_metadata()
    metadata.append(meta_entry)
    _save_metadata(metadata)

    logger.info('템플릿 저장 완료: %s (%s) — %d장, %d슬라이드',
                filename, template_id, analysis['chapter_count'], analysis['slide_count'])
    return meta_entry


def list_templates() -> list[dict]:
    """저장된 템플릿 목록을 반환합니다."""
    return _load_metadata()


def get_template_detail(template_id: str) -> dict | None:
    """템플릿 상세정보(장별 슬라이드 구조 포함)를 반환합니다."""
    metadata = _load_metadata()
    meta_entry = None
    for entry in metadata:
        if entry['id'] == template_id:
            meta_entry = entry
            break

    if meta_entry is None:
        return None

    analysis_path = os.path.join(_template_dir(template_id), 'analysis.json')
    if not os.path.exists(analysis_path):
        return meta_entry

    with open(analysis_path, 'r', encoding='utf-8') as f:
        analysis = json.load(f)

    return {**meta_entry, 'analysis': analysis}


def get_template_pptx_bytes(template_id: str) -> bytes | None:
    """저장된 PPTX 파일 바이트를 반환합니다."""
    pptx_path = os.path.join(_template_dir(template_id), 'template.pptx')
    if not os.path.exists(pptx_path):
        return None
    with open(pptx_path, 'rb') as f:
        return f.read()


def get_slide_thumbnail(template_id: str, slide_num: int,
                        width: int = 480) -> bytes | None:
    """특정 슬라이드의 썸네일 PNG를 반환합니다. 캐시 지원."""
    tpl_dir = _template_dir(template_id)
    cache_dir = os.path.join(tpl_dir, 'thumbnails')
    os.makedirs(cache_dir, exist_ok=True)

    cache_path = os.path.join(cache_dir, f'slide_{slide_num}_{width}.png')

    # 캐시된 썸네일이 있으면 반환
    if os.path.exists(cache_path):
        with open(cache_path, 'rb') as f:
            return f.read()

    # PPTX 파일 로드 후 생성
    pptx_bytes = get_template_pptx_bytes(template_id)
    if pptx_bytes is None:
        return None

    png_bytes = generate_slide_thumbnail(pptx_bytes, slide_num, width)
    if png_bytes is None:
        return None

    # 캐시 저장
    with open(cache_path, 'wb') as f:
        f.write(png_bytes)

    return png_bytes


def delete_template(template_id: str) -> bool:
    """템플릿을 삭제합니다."""
    metadata = _load_metadata()
    new_metadata = [m for m in metadata if m['id'] != template_id]

    if len(new_metadata) == len(metadata):
        return False

    _save_metadata(new_metadata)

    # 템플릿 디렉터리 전체 삭제
    import shutil
    tpl_dir = os.path.join(TEMPLATES_STORAGE_DIR, template_id)
    if os.path.exists(tpl_dir):
        shutil.rmtree(tpl_dir, ignore_errors=True)

    # 구버전 호환: 단일 파일 삭제
    for suffix in ['.pptx', '_analysis.json']:
        old_path = os.path.join(TEMPLATES_STORAGE_DIR, f'{template_id}{suffix}')
        if os.path.exists(old_path):
            os.remove(old_path)

    logger.info('템플릿 삭제 완료: %s', template_id)
    return True
