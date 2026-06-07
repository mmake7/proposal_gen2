"""제안서 PPTX 생성 경로 스모크 테스트.

제품 최종 산출물(제안서 .pptx)을 만드는 services.pptx_filler.fill_template을
LLM/네트워크 없이 검증한다: 최소 템플릿 + 샘플 매핑/콘텐츠/TOC → 유효한
.pptx(zip 'PK' 시그니처, python-pptx로 재오픈 가능)인지 확인.

경로(TEMPLATES_STORAGE_DIR)와 폰트 설정은 tmp로 격리해 실제 data/ 를 건드리지 않는다.
"""

import io

import pytest
from pptx import Presentation

import services.pptx_filler as filler


def _make_template(dir_path) -> str:
    """제목 슬라이드 + 제목/본문 슬라이드로 구성된 최소 템플릿 .pptx 생성."""
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])   # Title Slide (title idx 0)
    prs.slides.add_slide(prs.slide_layouts[1])   # Title and Content (title 0 / body 1)
    out = dir_path / 'template.pptx'
    prs.save(str(out))
    return str(out)


def test_fill_template_produces_valid_pptx(tmp_path, monkeypatch):
    storage = tmp_path / 'templates_storage'
    tid = 'smoke'
    tdir = storage / tid
    tdir.mkdir(parents=True)
    _make_template(tdir)

    # 경로/폰트 격리 (실제 data/ 미사용, 프리셋 분기 회피)
    monkeypatch.setattr(filler, 'TEMPLATES_STORAGE_DIR', str(storage))
    monkeypatch.setattr(filler, 'get_font_settings', lambda: {'mode': 'template_keep'})

    mapping = [{
        'status': 'matched', 'toc_index': 0,
        'toc_title': 'I. 사업 이해', 'slides': [1, 2],
    }]
    toc_items = [{
        'level': 1, 'title': 'I. 사업 이해',
        'children': [{'level': 2, 'title': '사업 배경', 'children': []}],
    }]
    content = {'sections': {'0-0': {
        'title': '사업 배경',
        'content_text': '본 사업은 ... 추진한다.',
        'bullet_points': ['핵심 포인트 1', '핵심 포인트 2'],
        'key_phrases': ['스마트시티', '플랫폼'],
    }}}

    buf = filler.fill_template(tid, mapping, content, toc_items)
    data = buf.read()

    # zip(PK) 시그니처 + python-pptx 재오픈 가능
    assert data[:2] == b'PK', '유효한 .pptx(zip) 산출물이 아님'
    reopened = Presentation(io.BytesIO(data))
    assert len(reopened.slides) == 2

    # 제목 슬라이드에 챕터 제목이 들어갔는지 (텍스트 어딘가에 포함)
    title_slide_text = ' '.join(
        sh.text_frame.text for sh in reopened.slides[0].shapes if sh.has_text_frame
    )
    assert '사업 이해' in title_slide_text


def test_fill_template_missing_template_raises(tmp_path, monkeypatch):
    """템플릿 파일이 없으면 PptxFillerError."""
    monkeypatch.setattr(filler, 'TEMPLATES_STORAGE_DIR', str(tmp_path))
    monkeypatch.setattr(filler, 'get_font_settings', lambda: {'mode': 'template_keep'})
    with pytest.raises(filler.PptxFillerError):
        filler.fill_template('does-not-exist', [], {'sections': {}}, [])


def test_fill_template_skips_unmatched_mapping(tmp_path, monkeypatch):
    """unmatched 매핑은 건너뛰고도 유효한 .pptx를 생성한다."""
    storage = tmp_path / 'templates_storage'
    tdir = storage / 'smoke'
    tdir.mkdir(parents=True)
    _make_template(tdir)
    monkeypatch.setattr(filler, 'TEMPLATES_STORAGE_DIR', str(storage))
    monkeypatch.setattr(filler, 'get_font_settings', lambda: {'mode': 'template_keep'})

    mapping = [{'status': 'unmatched', 'toc_index': 0, 'toc_title': 'x', 'slides': [1]}]
    buf = filler.fill_template('smoke', mapping, {'sections': {}}, [])
    assert buf.read()[:2] == b'PK'


def test_placeholder_classification_disjoint_and_correct():
    """제목/본문 placeholder가 타입 기준으로 상호배타 분류되고 텍스트가 올바른 자리에 들어간다.

    (과거 idx 기반: idx 1을 제목·본문 양쪽으로 보아 부제를 제목으로 오인하는 버그)
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content

    filler._fill_content_slide(
        slide,
        {'title': '섹션제목', 'content_text': '본문문장', 'bullet_points': ['불릿A']},
        'template_keep', None,
    )

    title_ph = body_ph = None
    for sh in slide.shapes:
        if not sh.is_placeholder:
            continue
        is_t = filler._is_title_placeholder(sh)
        is_b = filler._is_body_placeholder(sh)
        assert not (is_t and is_b), '제목/본문 분류가 겹침'
        if is_t:
            title_ph = sh
        elif is_b:
            body_ph = sh

    assert title_ph is not None and '섹션제목' in title_ph.text_frame.text
    assert body_ph is not None
    assert '본문문장' in body_ph.text_frame.text
    assert '불릿A' in body_ph.text_frame.text
